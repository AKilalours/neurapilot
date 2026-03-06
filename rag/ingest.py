"""Document ingestion pipeline for NeuraPilot.

Pipeline:
  1. Discover files (PDF, Markdown, TXT) in course upload directory
  2. Load documents using appropriate loaders
  3. Chunk with RecursiveCharacterTextSplitter
  4. Attach metadata: source, page, chunk_id, course_id, ingestion_ts
  5. Upsert into ChromaDB (idempotent via chunk_id as document ID)

Design notes:
- Idempotent: re-ingesting the same files produces the same chunk IDs
  (content-hashed), so no duplicate vectors are created.
- Progress: uses tqdm for CLI feedback during ingestion.
- Extensible: new file types can be added to _LOADERS without touching other code.
"""
from __future__ import annotations

import hashlib
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from tqdm import tqdm

from neurapilot.config import Settings
from neurapilot.rag.llm import LLMBundle
from neurapilot.rag.store import delete_course_collection, get_vector_store
from neurapilot.storage import course_upload_dir


# ── Supported file types ──────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS: frozenset[str] = frozenset({
    # Text / Markdown
    ".pdf", ".txt", ".md", ".rst",
    # Spreadsheets / Data
    ".csv", ".xlsx", ".xls",
    # Presentations
    ".pptx", ".ppt",
    # Word docs
    ".docx", ".doc",
    # Structured data
    ".json", ".jsonl",
    ".xml",
    ".yaml", ".yml",
    # Code files (treated as text)
    ".py", ".r", ".js", ".ts", ".html", ".htm", ".sql", ".ipynb",
})


@dataclass(frozen=True)
class IngestStats:
    """Summary of an ingestion run."""
    files_seen: int
    docs_loaded: int
    chunks_created: int
    chunks_indexed: int
    skipped_files: int
    duration_s: float

    def __str__(self) -> str:
        return (
            f"IngestStats(files={self.files_seen}, docs={self.docs_loaded}, "
            f"chunks={self.chunks_created}, indexed={self.chunks_indexed}, "
            f"skipped={self.skipped_files}, duration={self.duration_s:.1f}s)"
        )


# ── Loaders ───────────────────────────────────────────────────────────────────


def _load_pdf(path: Path) -> list[Document]:
    from langchain_community.document_loaders import PyPDFLoader
    return PyPDFLoader(str(path)).load()


def _load_text(path: Path) -> list[Document]:
    from langchain_community.document_loaders import TextLoader
    return TextLoader(str(path), encoding="utf-8", autodetect_encoding=True).load()


def _load_csv(path: Path) -> list[Document]:
    """Load CSV intelligently — batches rows, caps large datasets, adds stats summary.

    Strategy:
    - Small CSV (<500 rows): load all, batch 50 rows per document
    - Large CSV (≥500 rows): load first 200 + random sample 200 + column stats
    - This keeps ingestion fast while preserving searchable content
    """
    import csv
    import random

    MAX_ROWS    = 400   # total rows to index
    BATCH_SIZE  = 50    # rows per document chunk
    SAMPLE_SEED = 42

    try:
        # First pass: read all rows (fast, no embeddings yet)
        with open(path, encoding="utf-8", errors="replace") as f:
            reader  = csv.DictReader(f)
            headers = list(reader.fieldnames or [])
            all_rows = list(reader)

        total_rows  = len(all_rows)
        header_str  = " | ".join(headers)

        # Select rows to index
        if total_rows <= MAX_ROWS:
            selected = all_rows
            sampling_note = f"All {total_rows} rows indexed."
        else:
            # First 200 rows + random sample of 200
            head_rows   = all_rows[:200]
            remaining   = all_rows[200:]
            random.seed(SAMPLE_SEED)
            sample_size = min(200, len(remaining))
            sampled     = random.sample(remaining, sample_size)
            selected    = head_rows + sampled
            sampling_note = (
                f"Dataset has {total_rows} rows. "
                f"Indexed first 200 + {sample_size} random samples = {len(selected)} rows total."
            )

        # Compute column stats for numeric columns
        stats_lines: list[str] = []
        for col in headers:
            vals = []
            for row in all_rows:
                v = row.get(col, "")
                try:
                    vals.append(float(v))
                except (ValueError, TypeError):
                    pass
            if len(vals) > 5:
                mn  = min(vals)
                mx  = max(vals)
                avg = sum(vals) / len(vals)
                stats_lines.append(f"  {col}: min={mn:.2f}, max={mx:.2f}, avg={avg:.2f}, n={len(vals)}")

        # Build documents — one per batch of BATCH_SIZE rows
        docs: list[Document] = []

        # Document 0: schema + stats + sampling info
        schema_parts = [
            f"File: {path.name}",
            f"Total rows in file: {total_rows}",
            f"Columns ({len(headers)}): {header_str}",
            sampling_note,
        ]
        if stats_lines:
            schema_parts.append("Numeric column statistics:\n" + "\n".join(stats_lines))

        docs.append(Document(
            page_content="\n".join(schema_parts),
            metadata={"source": str(path), "filename": path.name, "type": "csv", "page": 0},
        ))

        # Documents 1..N: row batches
        for batch_start in range(0, len(selected), BATCH_SIZE):
            batch = selected[batch_start : batch_start + BATCH_SIZE]
            batch_end = batch_start + len(batch)
            row_lines: list[str] = []
            for i, row in enumerate(batch, start=batch_start + 1):
                cols = " | ".join(f"{k}: {v}" for k, v in row.items() if v and str(v).strip())
                if cols:
                    row_lines.append(f"Row {i}: {cols}")

            if row_lines:
                content = (
                    f"File: {path.name} | Rows {batch_start+1}-{batch_end}\n"
                    f"Columns: {header_str}\n\n"
                    + "\n".join(row_lines)
                )
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": str(path), "filename": path.name,
                        "type": "csv", "page": len(docs),
                        "row_start": batch_start + 1, "row_end": batch_end,
                    },
                ))

        print(f"[INFO] CSV {path.name}: {total_rows} rows → {len(docs)} docs ({sampling_note})")
        return docs

    except Exception as e:
        print(f"[WARN] CSV parse error {path.name}: {e}")
        return []


def _load_xlsx(path: Path) -> list[Document]:
    """Load Excel — each sheet batched, large sheets sampled."""
    MAX_ROWS   = 400
    BATCH_SIZE = 50

    try:
        import openpyxl
        wb   = openpyxl.load_workbook(path, read_only=True, data_only=True)
        docs: list[Document] = []

        for sheet_name in wb.sheetnames:
            ws      = wb[sheet_name]
            headers: list[str] = []
            all_rows: list[list[str]] = []

            for i, row in enumerate(ws.iter_rows(values_only=True)):
                vals = [str(c) if c is not None else "" for c in row]
                if i == 0:
                    headers = vals
                else:
                    all_rows.append(vals)

            total       = len(all_rows)
            selected    = all_rows[:MAX_ROWS]
            header_str  = " | ".join(headers)
            sample_note = f"All {total} rows indexed." if total <= MAX_ROWS else f"{total} rows total, first {MAX_ROWS} indexed."

            # Schema doc
            docs.append(Document(
                page_content=f"File: {path.name} | Sheet: {sheet_name}\nColumns: {header_str}\n{sample_note}",
                metadata={"source": str(path), "filename": path.name, "sheet": sheet_name, "type": "xlsx", "page": 0},
            ))

            # Row batch docs
            for b in range(0, len(selected), BATCH_SIZE):
                batch    = selected[b : b + BATCH_SIZE]
                lines    = []
                for j, row in enumerate(batch, start=b+1):
                    cols = " | ".join(f"{h}: {v}" for h, v in zip(headers, row) if v.strip())
                    if cols:
                        lines.append(f"Row {j}: {cols}")
                if lines:
                    docs.append(Document(
                        page_content=f"File: {path.name} | Sheet: {sheet_name} | Rows {b+1}-{b+len(batch)}\nColumns: {header_str}\n\n" + "\n".join(lines),
                        metadata={"source": str(path), "filename": path.name, "sheet": sheet_name, "type": "xlsx", "page": len(docs)},
                    ))

        return docs
    except Exception as e:
        print(f"[WARN] XLSX parse error {path.name}: {e}")
        return []


def _load_pptx(path: Path) -> list[Document]:
    """Load PowerPoint — each slide becomes a document with its text."""
    try:
        from pptx import Presentation  # type: ignore
        prs  = Presentation(str(path))
        docs: list[Document] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                content = f"File: {path.name} | Slide {i}\n\n" + "\n".join(texts)
                docs.append(Document(
                    page_content=content,
                    metadata={"source": str(path), "filename": path.name, "page": i, "type": "pptx"},
                ))
        return docs
    except Exception as e:
        print(f"[WARN] PPTX parse error {path.name}: {e}")
        return []


def _load_docx(path: Path) -> list[Document]:
    """Load Word document — extract all paragraphs and tables."""
    try:
        import docx  # type: ignore
        doc    = docx.Document(str(path))
        blocks: list[str] = []
        # Paragraphs
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                blocks.append(t)
        # Tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    blocks.append(row_text)
        content = f"File: {path.name}\n\n" + "\n\n".join(blocks)
        return [Document(page_content=content, metadata={"source": str(path), "filename": path.name, "type": "docx"})]
    except Exception as e:
        print(f"[WARN] DOCX parse error {path.name}: {e}")
        return []


_LOADERS: dict[str, Callable[[Path], list[Document]]] = {
    ".pdf":  _load_pdf,
    ".txt":  _load_text,
    ".md":   _load_text,
    ".rst":  _load_text,
    ".csv":  _load_csv,
    ".xlsx": _load_xlsx,
    ".xls":  _load_xlsx,
    ".pptx": _load_pptx,
    ".ppt":  _load_pptx,
    ".docx": _load_docx,
    ".doc":  _load_docx,
}

# ── Additional loaders ────────────────────────────────────────────────────────


def _load_json(path: Path) -> list[Document]:
    """Load JSON/JSONL — pretty-prints structure for semantic search."""
    import json as _json
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        # JSONL (one object per line)
        if path.suffix.lower() == ".jsonl":
            lines = [l.strip() for l in raw.splitlines() if l.strip()]
            docs: list[Document] = []
            for i, line in enumerate(lines[:200]):      # cap at 200 records
                try:
                    obj = _json.loads(line)
                    docs.append(Document(
                        page_content=f"File: {path.name} | Record {i+1}\n{_json.dumps(obj, indent=2)[:2000]}",
                        metadata={"source": str(path), "filename": path.name, "type": "jsonl", "page": i},
                    ))
                except Exception:
                    pass
            return docs
        else:
            obj = _json.loads(raw)
            pretty = _json.dumps(obj, indent=2)
            # Split large JSON into chunks of 3000 chars
            chunks: list[Document] = []
            for i in range(0, len(pretty), 3000):
                chunks.append(Document(
                    page_content=f"File: {path.name} | Part {i//3000+1}\n{pretty[i:i+3000]}",
                    metadata={"source": str(path), "filename": path.name, "type": "json", "page": i//3000},
                ))
            return chunks
    except Exception as e:
        print(f"[WARN] JSON parse error {path.name}: {e}")
        return _load_text(path)      # fallback to raw text


def _load_yaml(path: Path) -> list[Document]:
    """Load YAML — converts to readable key:value text."""
    try:
        import yaml  # type: ignore
        raw = path.read_text(encoding="utf-8", errors="replace")
        obj = yaml.safe_load(raw)
        import json as _json
        readable = _json.dumps(obj, indent=2, default=str)
        return [Document(
            page_content=f"File: {path.name}\n\n{readable[:8000]}",
            metadata={"source": str(path), "filename": path.name, "type": "yaml"},
        )]
    except Exception:
        return _load_text(path)


def _load_xml(path: Path) -> list[Document]:
    """Load XML — extracts text content from tags."""
    try:
        import xml.etree.ElementTree as ET
        tree  = ET.parse(str(path))
        root  = tree.getroot()
        texts: list[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                texts.append(f"{elem.tag}: {elem.text.strip()}")
            for k, v in elem.attrib.items():
                texts.append(f"{elem.tag}@{k}: {v}")
        content = f"File: {path.name}\n\n" + "\n".join(texts[:500])
        return [Document(page_content=content, metadata={"source": str(path), "filename": path.name, "type": "xml"})]
    except Exception:
        return _load_text(path)


def _load_notebook(path: Path) -> list[Document]:
    """Load Jupyter notebook — extracts code cells and markdown."""
    import json as _json
    try:
        nb = _json.loads(path.read_text(encoding="utf-8", errors="replace"))
        cells = nb.get("cells", [])
        docs: list[Document] = []
        for i, cell in enumerate(cells):
            ct    = cell.get("cell_type","")
            src   = "".join(cell.get("source",[]))
            if not src.strip(): continue
            label = "💻 Code" if ct == "code" else "📝 Markdown"
            docs.append(Document(
                page_content=f"File: {path.name} | Cell {i+1} ({label})\n{src[:3000]}",
                metadata={"source": str(path), "filename": path.name, "type": "ipynb", "page": i},
            ))
        return docs
    except Exception:
        return _load_text(path)


# Register additional loaders
_LOADERS.update({
    ".json":  _load_json,
    ".jsonl": _load_json,
    ".yaml":  _load_yaml,
    ".yml":   _load_yaml,
    ".xml":   _load_xml,
    ".ipynb": _load_notebook,
    # Code files → plain text
    ".py":    _load_text,
    ".r":     _load_text,
    ".js":    _load_text,
    ".ts":    _load_text,
    ".html":  _load_text,
    ".htm":   _load_text,
    ".sql":   _load_text,
})


def _load_file(path: Path) -> list[Document]:
    """Load a single file; returns empty list on error (with warning)."""
    ext = path.suffix.lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        return []
    try:
        docs = loader(path)
        for doc in docs:
            doc.metadata.setdefault("source", str(path))
            doc.metadata.setdefault("filename", path.name)
        return docs
    except Exception as exc:
        print(f"[WARN] Failed to load {path.name}: {exc}")
        return []


# ── Chunk ID generation ───────────────────────────────────────────────────────


def _stable_chunk_id(course_id: str, source: str, index: int, content: str) -> str:
    """Generate a stable, content-based chunk ID for idempotent upserts.

    The hash includes course_id + source + chunk content so that re-ingesting
    unchanged files produces the same IDs and avoids duplicate vectors.
    """
    payload = f"{course_id}|{source}|{index}|{content[:200]}"
    return hashlib.sha1(payload.encode()).hexdigest()[:16]


# ── Main ingestion function ───────────────────────────────────────────────────


def ingest_course(
    course_id: str,
    settings: Settings,
    bundle: LLMBundle,
    clear_existing: bool = False,
) -> IngestStats:
    """Run full ingestion pipeline for a course.

    Fast path for large files:
    - CSV/XLSX: row-batched with sampling (see loaders)
    - All formats: hard cap of 500 chunks to prevent runaway embedding time
    """
    t0 = time.time()

    if clear_existing:
        delete_course_collection(settings, bundle.embeddings, course_id)

    store = get_vector_store(settings, bundle.embeddings, course_id)
    root  = course_upload_dir(settings, course_id)
    files = [
        p for p in root.rglob("*")
        if p.is_file()
        and p.suffix.lower() in SUPPORTED_EXTENSIONS
        and not p.name.startswith(".")
    ]

    all_docs: list[Document] = []
    skipped = 0

    for fp in tqdm(files, desc=f"Loading [{course_id}]", unit="file"):
        loaded = _load_file(fp)
        if loaded:
            all_docs.extend(loaded)
            print(f"[INFO] {fp.name}: {len(loaded)} doc(s) loaded")
        else:
            skipped += 1

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        add_start_index=True,
        separators=["\n\n", "\n", ". ", "? ", "! ", " ", ""],
    )
    chunks = splitter.split_documents(all_docs)

    # Safety cap: never embed more than 600 chunks (prevents multi-hour ingestion)
    # For CSV/XLSX this is already handled by sampling in the loaders
    MAX_CHUNKS = 600
    if len(chunks) > MAX_CHUNKS:
        print(f"[INFO] Capping {len(chunks)} chunks → {MAX_CHUNKS} (change MAX_CHUNKS in ingest.py if needed)")
        chunks = chunks[:MAX_CHUNKS]

    # Attach stable IDs and metadata
    ts = int(time.time())
    for idx, chunk in enumerate(chunks):
        source   = str(chunk.metadata.get("source", "unknown"))
        chunk_id = _stable_chunk_id(course_id, source, idx, chunk.page_content)
        chunk.metadata["chunk_id"]      = chunk_id
        chunk.metadata["course_id"]     = course_id
        chunk.metadata["ingestion_ts"]  = ts
        chunk.metadata["chunk_index"]   = idx

    ids = [c.metadata["chunk_id"] for c in chunks]
    if chunks:
        # Batch upsert in groups of 100 to show progress
        BATCH = 100
        for i in range(0, len(chunks), BATCH):
            batch = chunks[i:i+BATCH]
            store.add_documents(batch, ids=[c.metadata["chunk_id"] for c in batch])
            print(f"[INFO] Embedded {min(i+BATCH, len(chunks))}/{len(chunks)} chunks…")

    return IngestStats(
        files_seen=len(files),
        docs_loaded=len(all_docs),
        chunks_created=len(chunks),
        chunks_indexed=len(chunks),
        skipped_files=skipped,
        duration_s=round(time.time() - t0, 2),
    )
